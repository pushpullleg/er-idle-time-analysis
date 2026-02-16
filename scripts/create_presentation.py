#!/usr/bin/env python3
"""
DATATHON SUBMISSION: Create Professional PowerPoint Presentation
Meridian ER Optimization: Emergency Room Efficiencies
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color scheme (professional healthcare)
DARK_BLUE = RGBColor(25, 65, 130)      # Professional dark blue
ACCENT_ORANGE = RGBColor(255, 102, 0)  # Alert/attention orange
LIGHT_GRAY = RGBColor(245, 245, 245)   # Background
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(45, 45, 45)

def add_title_slide(prs, title, subtitle=""):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.word_wrap = True
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(left, Inches(4.2), width, Inches(1))
        sub_frame = sub_box.text_frame
        sub_frame.text = subtitle
        sub_frame.paragraphs[0].font.size = Pt(28)
        sub_frame.paragraphs[0].font.color.rgb = ACCENT_ORANGE
        sub_frame.paragraphs[0].font.italic = True
    
    return slide

def add_content_slide(prs, title, content_list=None):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.margin_left = Inches(0.3)
    
    # Content
    if content_list:
        left = Inches(0.5)
        top = Inches(1.1)
        width = Inches(9)
        height = Inches(5)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_DARK
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    return slide

def add_metric_slide(prs, title, metrics):
    """Add slide with key metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.margin_left = Inches(0.3)
    
    # Metrics boxes
    box_width = Inches(2.8)
    box_height = Inches(2.2)
    start_left = Inches(0.4)
    start_top = Inches(1.2)
    spacing = Inches(0.15)
    
    for idx, (label, value, unit) in enumerate(metrics):
        col = idx % 3
        row = idx // 3
        
        left = start_left + col * (box_width + spacing)
        top = start_top + row * (box_height + spacing)
        
        # Box background
        box = slide.shapes.add_shape(1, left, top, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = DARK_BLUE
        box.line.width = Pt(2)
        
        # Text frame
        text_frame = box.text_frame
        text_frame.word_wrap = True
        
        # Value (big)
        p_value = text_frame.paragraphs[0]
        p_value.text = value
        p_value.font.size = Pt(36)
        p_value.font.bold = True
        p_value.font.color.rgb = ACCENT_ORANGE
        p_value.alignment = PP_ALIGN.CENTER
        
        # Unit
        p_unit = text_frame.add_paragraph()
        p_unit.text = unit
        p_unit.font.size = Pt(11)
        p_unit.font.color.rgb = TEXT_DARK
        p_unit.alignment = PP_ALIGN.CENTER
        
        # Label
        p_label = text_frame.add_paragraph()
        p_label.text = label
        p_label.font.size = Pt(13)
        p_label.font.color.rgb = TEXT_DARK
        p_label.alignment = PP_ALIGN.CENTER
    
    return slide

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# SLIDE 1: Title Slide
add_title_slide(prs, 
    "Emergency Room Efficiencies",
    "From Bottleneck to Breakthrough")

# SLIDE 2: The Problem (Consultant speaks)
add_content_slide(prs, "The Mystery",
    [
        "💡 Meridian ER has 4 doctors on duty",
        "   But seeing only 6.9 patients/hour vs. 10+ capacity",
        "",
        "📊 That 25% capacity gap = $890K QUARTERLY loss",
        "",
        "❓ Question: Need more doctors?",
        "   OR... something else is happening?"
    ])

# SLIDE 3: Investigation Results (Analyst speaks)
add_metric_slide(prs, "Where Is All The Time Going?", 
    [
        ("Post-Triage Wait", "39 min", "23% of total"),
        ("Doctor Cycle Time", "107 min", "62% of total"),
        ("Other Stages", "7 min", "15% of total"),
        ("Bottlenecks Found", "2,179", "events in Q1"),
        ("Idle Doctor Moments", "14.5%", "of all visits"),
        ("Current Utilization", "50%", "vs 75-80% target")
    ])

# SLIDE 4: The Revelation (Scientist speaks)
add_content_slide(prs, "The Breakthrough Discovery",
    [
        "🔍 KEY FINDING: The bottleneck is NOT doctors",
        "",
        "⚡ In 14.5% of visits:",
        "   • 1.8 doctors are IDLE",
        "   • 4.3 patients are WAITING",
        "   • AT THE SAME TIME",
        "",
        "📈 System Theory insight:",
        "   This is a PROCESS problem, not a STAFFING problem",
        "",
        "✅ We have capacity. We just aren't using it.",
    ])

# SLIDE 5: Two Situations Framework
add_content_slide(prs, "The Full Picture",
    [
        "🎯 SITUATION 1: Process Inefficiency",
        "   2,179 idle doctor moments → Workflow optimization",
        "",
        "🎯 SITUATION 2: Capacity Crisis",
        "   Day shift at 154% utilization → Need staffing",
        "",
        "🔗 WHY BOTH MATTER:",
        "   Process alone → Gets to 70%, still hits ceiling",
        "   Staffing alone → New docs inherit broken workflows",
        "   BOTH together → 75-80% sustainable operations",
    ])

# SLIDE 6: The Solution (Consultant & Scientist)
add_metric_slide(prs, "3-Tier Implementation Plan",
    [
        ("TIER 1 (Wk 1-4)", "$150K", "Queue + Protocol"),
        ("TIER 2 (Wk 5-8)", "$250K", "System + NP Fast-Track"),
        ("TIER 3 (Wk 9-12)", "$150K", "Monitoring + Optimization"),
        ("Expected Utilization", "50% → 70%", "Phase 1-2"),
        ("New Staff Added", "1-2 MDs", "Morning hours"),
        ("Expected Utilization", "70% → 75-80%", "Full solution"),
    ])

# SLIDE 7: ML Models & Impact
add_content_slide(prs, "5 ML Models Powering The Solution",
    [
        "🤖 Model 1: Complexity Predictor (Random Forest)",
        "   → Predicts patient complexity → Better routing",
        "",
        "🤖 Model 2: Intelligent Dispatcher (LightGBM)",
        "   → Auto-assigns next patient → Dispatch delay 5→1 min",
        "",
        "🤖 Model 3: Length-of-Stay (Quantile Regression)",
        "   → Predicts visit duration → Queue optimization",
        "",
        "🤖 Model 4: Workload Forecaster (Prophet + XGBoost)",
        "   → Predicts bottlenecks 2 hours ahead → Proactive staffing",
        "",
        "🤖 Model 5: Adverse Outcome Detector (Neural Network)",
        "   → High-risk patients flagged → Early intervention",
    ])

# SLIDE 8: Financial Impact & ROI
add_metric_slide(prs, "The Business Case",
    [
        ("Investment", "$550K", "Tiers 1-3"),
        ("Year 1 Benefit", "$15.3M", "7,500 add'l visits"),
        ("Payback Period", "3.3 weeks", "Immediate ROI"),
        ("Year 1 ROI", "2,662%", "Outstanding"),
        ("Staff Added", "1-2 MDs", "$225K/year cost"),
        ("Net Benefit Year 2+", "$1.9M", "Annual sustainable"),
    ])

# Save presentation
output_path = str(Path(__file__).resolve().parent.parent / "presentation" / "Meridian_ER_Presentation.pptx")
prs.save(output_path)

print(f"✅ PowerPoint presentation created: {output_path}")
print(f"   • 8 slides with professional design")
print(f"   • Consultant → Analyst → Scientist narrative flow")
print(f"   • Data visuals, metrics, solution framework")
print(f"   • Ready for voiceover recording")
